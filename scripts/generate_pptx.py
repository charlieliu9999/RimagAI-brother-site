#!/usr/bin/env python3
"""Generate RimagAi Brother hospital pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
PRIMARY = RGBColor(0x63, 0x66, 0xF1)  # Indigo
DARK = RGBColor(0x1E, 0x1B, 0x4B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)  # Cyan
GRAY = RGBColor(0x64, 0x74, 0x8B)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(6)
    return txBox


def add_card(slide, left, top, width, height, title, body, title_color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(12)
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Microsoft YaHei"

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = "Microsoft YaHei"
    p2.space_before = Pt(6)


prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "RimagAi Brother", font_size=44, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.0),
             "医疗认知操作系统 MCOS", font_size=28, bold=False, color=ACCENT,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.8),
             "告别单点 AI 插件，构建系统级医疗人工智能", font_size=18, color=WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
             "院内 AI 落地方案 · 2026", font_size=14, color=GRAY,
             alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "医院 AI 的三大困境", font_size=32, bold=True, color=DARK)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             '核心判断：医疗 AI 的主要矛盾已从"模型能力不足"转变为"架构与治理能力缺位"',
             font_size=14, color=GRAY)

problems = [
    ("工具孤岛", "各科室独立采购 AI 工具，彼此不通\n数据无法流转，经验无法共享"),
    ("治理真空", "AI 生成内容无统一审计\n无法追溯、无法问责、无法合规"),
    ("能力不可积累", "医生反馈无法沉淀为组织级知识\n系统越用越旧，而非越用越懂"),
]

for i, (title, body) in enumerate(problems):
    left = Inches(0.8 + i * 4.0)
    add_card(slide, left, Inches(2.0), Inches(3.6), Inches(3.5), title, body)

add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1.0),
             "→ 医院需要的不是更多 AI 工具，而是一个统一的智能·认知底座",
             font_size=16, bold=True, color=PRIMARY)

# ============================================================
# Slide 3: MCOS Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "MCOS 六层架构", font_size=32, bold=True, color=DARK)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "医院现有系统 (HIS/EMR/PACS/RIS/LIS) → MCOS 认知运行层 → AI 能力层",
             font_size=14, color=GRAY)

layers = [
    ("L6 进化层", "经验萃取、反馈学习、技能优化", "目标态"),
    ("L5 治理层", "权限分级、审计链、风险控制", "建设中"),
    ("L4 执行层", "建议、草稿、预填、受控写回", "已验证"),
    ("L3 认知层", "意图识别、推理、规划、工具调用", "已验证"),
    ("L2 上下文层", "患者纵向/当前/场景工作流上下文", "已验证"),
    ("L1 感知层", "多模态信号采集与预处理", "已验证"),
]

for i, (name, desc, status) in enumerate(layers):
    top = Inches(1.8 + i * 0.85)
    color = SUCCESS if status == "已验证" else (ACCENT if status == "建设中" else GRAY)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.0), top, Inches(11.0), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(16)
    p = tf.paragraphs[0]
    p.text = f"{name}  —  {desc}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = "Microsoft YaHei"

    add_text_box(slide, Inches(10.5), top + Emu(Inches(0.1).emu),
                 Inches(1.5), Inches(0.5), f"[{status}]",
                 font_size=11, color=color, alignment=PP_ALIGN.RIGHT)

# ============================================================
# Slide 4: Product Matrix
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "8 大核心产品 · 3 科室全覆盖", font_size=32, bold=True, color=DARK)

departments = [
    ("临床科室", [
        "CDSS 决策支持 — 分阶段推理，显式不确定性",
        "智能病历录入 — 语音+截屏+文字三通道并行",
        "DRG 费用预警 — 医嘱前置拦截，实时合规检查",
    ]),
    ("影像科室", [
        "智能报告生成 — 语音控制+DICOM集成，效率提升60%+",
        "双重质控 — 规则引擎+大模型语义质控",
        "ACRAC 检查推荐 — 基于ACR准则，覆盖200+场景",
    ]),
    ("信息科室", [
        "院内数据智能体 — 自然语言查询+权限矩阵",
        "病历内涵质控 — 后台自动轮询+问题分级推送",
    ]),
]

for i, (dept, items) in enumerate(departments):
    left = Inches(0.8 + i * 4.0)
    add_text_box(slide, left, Inches(1.5), Inches(3.6), Inches(0.5),
                 dept, font_size=18, bold=True, color=PRIMARY)
    add_bullet_list(slide, left, Inches(2.1), Inches(3.8), Inches(4.5),
                    [f"• {item}" for item in items], font_size=12)

# ============================================================
# Slide 5: Clinical CDSS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "CDSS 智能临床决策支持", font_size=32, bold=True, color=DARK)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "分阶段推理 · 显式不确定性 · 依据关联", font_size=16, color=GRAY)

stages = [
    ("① 诊断推理", "生成鉴别诊断列表\n置信度排序"),
    ("② 检查建议", "推荐必要检查项目\nACR适宜性等级"),
    ("③ 用药方案", "结合诊断推荐用药\n禁忌+相互作用检查"),
    ("④ 处置计划", "综合处置建议\n含随访计划"),
]

for i, (title, body) in enumerate(stages):
    left = Inches(0.8 + i * 3.1)
    add_card(slide, left, Inches(2.2), Inches(2.8), Inches(2.5), title, body)

features = [
    "• 每个推荐附带置信度评分",
    "• 推荐结果关联到具体指南/文献/本院规范",
    "• 规则引擎兜底 — 确定性逻辑零延迟执行",
    "• 上下文感知 — 自动整合患者纵向病史和HIS数据",
]
add_bullet_list(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(2.0),
                features, font_size=13)

# ============================================================
# Slide 6: Radiology Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "影像科方案 — 重塑报告工作流", font_size=32, bold=True, color=DARK)

cards_data = [
    ("智能语音报告生成", "• 医学专有名词增强（FunASR优化）\n• 在线/离线双模切换\n• 报告编写时间缩短 60%+\n• MCOS: L1感知层 + L4执行层"),
    ("智能双重防呆质控", "• 规则引擎毫秒级拦截（0.1s）\n• 大模型上下文一致性检查\n• 双重保障，互为补充\n• MCOS: L3认知层 + L5治理层"),
    ("文本知识结构化抽取", "• 多维度实体抽取\n  （病灶部位/性质/大小/分级）\n• 赋能临床科研闭环\n• MCOS: L2上下文层 + L6进化层"),
]

for i, (title, body) in enumerate(cards_data):
    left = Inches(0.8 + i * 4.0)
    add_card(slide, left, Inches(1.8), Inches(3.7), Inches(4.5), title, body)

# Key metrics
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
             ">98% 名词识别准确率  |  0.1s 规则拦截延迟  |  JSON 标准化资产产出",
             font_size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 7: Data Platform
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "院内数据智能体", font_size=32, bold=True, color=DARK)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "自然语言输入 → 权限校验 → 语义理解 → 受控执行 → 结果呈现",
             font_size=14, color=GRAY)

left_items = [
    "• 自然语言查询 — 支持中文自然语言提问",
    "• MCP 工具链 — 基于 MCP 协议的工具编排",
    "• Schema 感知 — 自动获取数据库表结构",
    "• 智能可视化 — 查询结果自动选择最佳展示方式",
]
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.0),
                left_items, font_size=13)

right_items = [
    "安全与治理：",
    "• 权限矩阵 — 基于角色的表/字段级访问控制",
    "• 审计日志 — 完整记录每次查询",
    "• SQL 安全沙箱 — 禁止 DELETE/UPDATE/DROP",
    "",
    "典型场景：",
    "• 各科室门诊量对比",
    "• 抗生素使用量 TOP10 科室",
    "• ICU 平均住院日和床位利用率",
    "• 非计划再入院率超标病区",
]
add_bullet_list(slide, Inches(6.5), Inches(2.0), Inches(5.5), Inches(5.0),
                right_items, font_size=13)

# ============================================================
# Slide 8: Deployment & Security
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "部署方式与数据安全", font_size=32, bold=True, color=DARK)

deploy_cards = [
    ("院内私有化部署", "• 核心数据不出院墙\n• 支持国产算力（昇腾/海光）\n• Docker Compose 一键编排\n• 离线镜像打包"),
    ("混合云部署", "• 敏感数据本地处理\n• 非敏感能力云端加速\n• 灵活切换，按需配置\n• 双通道 AI 架构支撑"),
    ("非侵入式对接", "• 不改动原有系统架构\n• API 层对接 HIS/EMR/PACS\n• 渐进式接入，风险可控\n• 支持多厂商系统并存"),
]

for i, (title, body) in enumerate(deploy_cards):
    left = Inches(0.8 + i * 4.0)
    add_card(slide, left, Inches(1.5), Inches(3.7), Inches(3.8), title, body)

security_items = [
    "本地优先 — 核心数据不出院墙",
    "可控可审计 — 完整操作日志 + 权限分级",
    "无缝旧系统融合 — 非侵入式 API 对接",
    "资产可积累 — 知识沉淀为医院自有资产",
]
add_bullet_list(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5),
                [f"✓ {item}" for item in security_items], font_size=14, color=SUCCESS)

# ============================================================
# Slide 9: 90-Day Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "90 天落地路径", font_size=32, bold=True, color=DARK)

phases = [
    ("P1 · 第 1-30 天", "基础接入 + 单场景验证",
     "• 环境部署与网络对接\n• HIS/EMR API 联调\n• 选定 1 个科室试点\n• 单场景功能验证"),
    ("P2 · 第 31-60 天", "多场景扩展 + 治理层上线",
     "• 扩展至 3+ 场景\n• 权限矩阵配置\n• 审计日志上线\n• 规则引擎调优"),
    ("P3 · 第 61-90 天", "全场景覆盖 + 进化层启动",
     "• 8 大产品全面上线\n• 知识库初始化\n• 反馈学习机制启动\n• 运营数据看板交付"),
]

for i, (title, subtitle, body) in enumerate(phases):
    left = Inches(0.8 + i * 4.0)
    add_text_box(slide, left, Inches(1.5), Inches(3.6), Inches(0.5),
                 title, font_size=18, bold=True, color=PRIMARY)
    add_text_box(slide, left, Inches(2.0), Inches(3.6), Inches(0.4),
                 subtitle, font_size=13, color=GRAY)
    add_bullet_list(slide, left, Inches(2.5), Inches(3.6), Inches(4.0),
                    body.split("\n"), font_size=12)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             "渐进式落地，每阶段可独立验收，风险可控",
             font_size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 10: Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "与现有方案对比", font_size=32, bold=True, color=DARK)

# Table header
headers = ["维度", "单点工具", "传统CDSS", "通用Agent框架", "MCOS"]
col_widths = [Inches(2.0), Inches(2.2), Inches(2.2), Inches(2.8), Inches(3.0)]
table_left = Inches(0.6)
table_top = Inches(1.5)

rows_data = [
    ["场景覆盖", "单一", "有限", "通用但缺医疗适配", "8个核心临床场景"],
    ["治理能力", "无", "基础审计", "无医疗治理", "权限+审计+规则引擎"],
    ["数据安全", "依赖厂商", "本地但封闭", "云端为主", "本地优先+国产算力"],
    ["系统集成", "独立运行", "深度耦合HIS", "需大量定制", "非侵入式对接"],
    ["知识积累", "不可积累", "手动维护", "无组织级沉淀", "自动经验萃取+反馈学习"],
]

total_width = sum(w.emu for w in col_widths)
table = slide.shapes.add_table(len(rows_data) + 1, 5, table_left, table_top,
                               Emu(total_width), Inches(4.5)).table

for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Header row
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY

# Data rows
for r, row in enumerate(rows_data):
    for c, val in enumerate(row):
        cell = table.cell(r + 1, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.name = "Microsoft YaHei"
        if c == 4:
            p.font.bold = True
            p.font.color.rgb = PRIMARY
        else:
            p.font.color.rgb = DARK

# ============================================================
# Slide 11: Key Metrics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "关键指标", font_size=32, bold=True, color=DARK)

metrics = [
    ("8", "核心产品"),
    ("3", "覆盖科室"),
    ("90天", "落地周期"),
    ("60%+", "效率提升"),
    (">98%", "名词识别"),
    ("0.1s", "规则拦截"),
    ("200+", "ACR场景"),
    ("58", "测试用例"),
]

for i, (num, label) in enumerate(metrics):
    col = i % 4
    row = i // 4
    left = Inches(1.0 + col * 3.0)
    top = Inches(1.8 + row * 2.5)

    add_text_box(slide, left, top, Inches(2.5), Inches(1.0),
                 num, font_size=36, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(1.0), Inches(2.5), Inches(0.5),
                 label, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 12: Summary & CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
             "为什么选择 MCOS？", font_size=36, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

reasons = [
    "统一编排 — 医生只需面对一个系统入口",
    "内建治理 — 权限分级、审计链、规则引擎",
    "本地优先 — 核心数据不出院墙",
    "持续进化 — 系统越用越懂本院业务",
    "90天落地 — 渐进式部署，每阶段可独立验收",
]

for i, reason in enumerate(reasons):
    add_text_box(slide, Inches(2.5), Inches(2.8 + i * 0.7), Inches(8), Inches(0.6),
                 f"✓  {reason}", font_size=16, color=WHITE)

add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.8),
             "让我们一起，构建系统级医疗人工智能", font_size=20, bold=True, color=ACCENT,
             alignment=PP_ALIGN.CENTER)

# Save
output_path = "/Users/charlieliu/git_project_vscode/09_medical/demo-web/rimagai-brother-site/docs/RimagAi-Brother-院内AI落地方案.pptx"
prs.save(output_path)
print(f"PPT generated: {output_path}")
